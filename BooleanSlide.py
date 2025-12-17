# Local
import constants as con
from make_ppt import make_title_box, make_title_layout
from Slide import Slide
from slideInfoGathering import TrueFalseReader

# Third Party
import openpyxl as px

from pptx import presentation, Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

# Standard Library
from enum import IntEnum
from pathlib import Path

class BooleanSlide(Slide):
    '''
    Reworked True/False Slide, showing the predictions in a less terrible
    format then a massive table. Presentation matters!
    '''

    class SlideType(IntEnum):
        PLAYER_PREDS = 0
        RESULT_REVEAL = 1
        SCORE_REVEAL = 2

    def __init__(self,
        prs: presentation.Presentation,
        driver_cond: dict[str, tuple[int, bool]],
        player_to_truths: dict[str, list[str]],
        player_scores: dict[str, int],
        base_title: str):

        self.prs = prs
        blank_slide_layout = prs.slide_layouts[6]
        self.slides = [prs.slides.add_slide(blank_slide_layout)
                    for i in range(len(BooleanSlide.SlideType))]

        self.driver_cond = driver_cond
        self.player_to_truths = player_to_truths
        self.player_scores = player_scores
        self.base_title = base_title

    def make_boxes(self):
        ph = Inches(1.0)
        slide = self.slides[0]
        shape_tree = slide.shapes

        self.displayed_drivers = [
            {},
            {player : {} for player in self.player_to_truths},
            {player : {} for player in self.player_to_truths}
        ]
        for player, drivers in self.player_to_truths.items():
            self.displayed_drivers[0][player] = {driver:
                shape_tree.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     ph, ph, ph, ph)
                        for driver in drivers}

        # Only on last slide do we show checks
        self.checks = {player : {} for player in self.player_to_truths}
        # And Counts
        self.counts = [{}, {}]
        for player, drivers in self.player_to_truths.items():
            for driver, condition in self.driver_cond.items():
                is_true = condition[1]
                if is_true or driver in drivers:
                    shape_tree = self.slides[1].shapes
                    self.displayed_drivers[1][player][driver] = shape_tree.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph)

                    if driver not in self.counts[0]:
                        self.counts[0][driver] = shape_tree.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph)

                    shape_tree = self.slides[2].shapes
                    self.displayed_drivers[2][player][driver] = shape_tree.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph)

                    if driver not in self.counts[1]:
                        self.counts[1][driver] = shape_tree.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph)

                # Need ✓ and ✗
                self.checks[player][driver] = shape_tree.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph)

        self.player_label_boxes = [{} for i in range(len(BooleanSlide.SlideType))]
        for i, slide in enumerate(self.slides):
            shape_tree = slide.shapes
            for player in self.player_scores:
                self.player_label_boxes[i][player] = shape_tree.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph)

    def make_layout(self):
        prs_width, prs_height = self.prs.slide_width, self.prs.slide_height

        # Where the first player should start is 1 - available_h
        available_h = 18/25

        # Available width for the drivers
        available_w = 19/20

        widths = int((available_w) * prs_width / (2 * len(self.driver_cond) - 1))
        heights = int(prs_height * available_h / (2 * len(self.player_to_truths) - 1))
        heights = min(widths, heights)

        padding = 1/20
        player_widths = int((1 - 2 * padding) * (1 - available_w) * prs_width)
        player_width_offset = int(padding * (1 - available_w) * prs_width)

        height_offset = int((1 - available_h) * prs_height)
        width_offset = int((1 - available_w) * prs_width)

        for i in range(len(self.slides)):
            for k, player in enumerate(self.player_scores):
                player_box = self.player_label_boxes[i][player]

                player_rgb = con.PLAYERS_TO_COLOURS[player]
                outline_rgb = (int(0.8 * i) for i in player_rgb)
                player_box.fill.solid()
                player_box.fill.fore_color.rgb = RGBColor(*player_rgb)
                player_box.line.color.rgb = RGBColor(*outline_rgb)

                player_heights = height_offset + 2 * k * heights

                player_box.left = player_width_offset
                player_box.top = player_heights
                player_box.width = player_widths
                player_box.height = heights

                displayed_drivers = self.displayed_drivers[i][player]
                for j, driver in enumerate(self.driver_cond):
                    if driver in displayed_drivers:
                        box = displayed_drivers[driver]
                        box.left = width_offset + 2 * j * widths
                        box.top = player_heights
                        box.width = widths
                        box.height = heights

                        if driver in self.player_to_truths[player]:
                            driver_rgb = con.TEAMS_TO_COLOURS[con.DRIVERS_TO_TEAMS[driver]][1:]
                            outline_rgb = driver_rgb
                            box.fill.solid()
                            box.fill.fore_color.rgb = RGBColor.from_string(driver_rgb)
                            box.line.color.rgb = RGBColor.from_string(outline_rgb)
                        else:
                            outline_rgb = con.TEAMS_TO_COLOURS[con.DRIVERS_TO_TEAMS[driver]][1:]
                            box.fill.background()
                            box.line.color.rgb = RGBColor.from_string(outline_rgb)

                        if i == 1:
                            box = self.counts[0][driver]
                            box.left = width_offset + 2 * j * widths - widths
                            box.top = height_offset + 2 * len(self.player_scores) * heights
                            box.width = 3 * widths
                            box.height = heights

                            box.fill.background()
                            box.line.fill.background()
                        elif i == 2:
                            box = self.counts[1][driver]
                            box.left = width_offset + 2 * j * widths - widths
                            box.top = height_offset + 2 * len(self.player_scores) * heights
                            box.width = 3 * widths
                            box.height = heights

                            box.fill.background()
                            box.line.fill.background()

                    if i == 2:
                        box = self.checks[player][driver]
                        box.left = width_offset + 2 * j * widths
                        box.top = player_heights
                        box.width = widths
                        box.height = heights

                        box.fill.background()
                        box.line.fill.background()

    def make_content(self):
        for i in range(len(self.slides)):
            for player in self.player_scores:
                player_box = self.player_label_boxes[i][player]
                if i != 2:
                    text = player[0]
                else:
                    text = str(self.player_scores[player])
                text_frame = player_box.text_frame
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = text

                displayed_drivers = self.displayed_drivers[i][player]
                for j, driver in enumerate(self.driver_cond):
                    if driver in displayed_drivers:
                        box = displayed_drivers[driver]
                        driver_initials = "".join(name[0] for name in driver.split(" "))
                        text_frame = box.text_frame
                        p = text_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        run = p.add_run()
                        run.text = driver_initials
                        run.font.size = Pt(8)

                    if i == 2:
                        box = self.checks[player][driver]
                        if self.driver_cond[driver][1] and driver in self.player_to_truths[player]:
                            text = "✓"
                        elif self.driver_cond[driver][1] or driver in self.player_to_truths[player]:
                            text = "✗"
                        else:
                            text = ""
                        text_frame = box.text_frame
                        p = text_frame.paragraphs[0]
                        p.alignment = PP_ALIGN.CENTER
                        run = p.add_run()
                        run.text = text
                        run.font.size = Pt(18)
                        run.font.color.rgb = RGBColor(255, 255, 0)

        for counts in self.counts:
            for driver, box in counts.items():
                text_frame = box.text_frame
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                run.text = str(self.driver_cond[driver][0])
                run.font.size = Pt(10)
                run.font.color.rgb = RGBColor(255, 255, 255)

    def make_titles(self):
        for slide in self.slides:
            title_box = make_title_layout(self.prs, slide, 4/25)
            make_title_box(self.base_title, title_box)

if __name__ == "__main__":

    prs = Presentation()
    prs.slide_width = 12192 * 1000
    blank_slide_layout = prs.slide_layouts[6]

    spreadsheet_path = Path("C:\Projekty\Coding\Python\F1PredsPPT\F12025 Predictions Tracking.xlsx")
    wb = px.open(spreadsheet_path, data_only=True)

    dc_reader = TrueFalseReader(wb["Q2Elims"])
    dc_reader.gather_data()
    data = dc_reader.format_to_slide()

    print(data[0])

    slide = BooleanSlide(
        prs,
        data[0],
        data[1],
        data[2],
        "True/False")

    slide.make_slide()

    prs.save("examples/bool.pptx")
