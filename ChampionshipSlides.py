# Local
import constants2025 as con
from ChampionshipReader import ChampionshipReader2025
from make_ppt import make_title_box, make_title_layout
from Slide import Slide

# Third Party
import openpyxl as px

from pptx import presentation, Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches

# Standard Lib
from pathlib import Path

class ChampionshipSlides(Slide):
    '''
    Revised style of Championship Slides.
    Each Competitor is ran through in reverse championship order, with individual
    people's predictions and scores displayed for each.
    '''

    def __init__(
            self,
            prs: presentation.Presentation,
            competitor_order: list[str],
            player_preds: dict[str, dict[str, tuple[int, int]]],
            running_scores: dict[str, dict[str, int]],
            base_title: str
        ):
        self.prs = prs
        blank_slide_layout = prs.slide_layouts[6]
        self.slides = [prs.slides.add_slide(blank_slide_layout) for competitor in competitor_order]

        self.competitor_order = competitor_order
        self.player_preds = player_preds
        self.running_scores = running_scores
        self.base_title = base_title

    def make_boxes(self):
        ph = Inches(1.0)
        self.player_boxes = []
        for slide in self.slides:
            shape_tree = slide.shapes
            self.player_boxes.append({player: shape_tree.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph
                ) for player in self.player_preds})

    def make_layout(self):
        num_players = len(self.player_preds)
        k = num_players // 2

        # if numplayers is odd
        if num_players % 2:
            divs = 3 * k + 4
            top_ls = list(range(1, divs, 3))
            bottom_ls = [i + 0.5 for i in range(2, 3 * k + 2, 3)]
        else:
            divs = 3 * k + 1
            top_ls = [i for i in range(1, divs, 3)]
            bottom_ls = top_ls

        available_h = 18/25
        delta_h = available_h / 5

        top_h = int((1 - available_h) * self.prs.slide_height)
        bottom_h = int(top_h + 2 * delta_h * self.prs.slide_height)

        width = int(2 / divs * self.prs.slide_width)
        height = int(delta_h * self.prs.slide_height)

        for i in range(len(self.slides)):
            j = 0
            print(self.player_boxes)
            for player, box in self.player_boxes[i].items():
                if j < len(top_ls):
                    left = int(top_ls[j] / divs * self.prs.slide_width)
                    top = top_h
                else:
                    left = int(bottom_ls[j - len(top_ls)] / divs * self.prs.slide_width)
                    top = bottom_h

                player_rgb = con.PLAYERS_TO_COLOURS[player]
                outline_rgb = (int(0.8 * i) for i in player_rgb)
                box.fill.solid()
                box.fill.fore_color.rgb = RGBColor(*player_rgb)
                box.line.color.rgb = RGBColor(*outline_rgb)

                box.left = left
                box.width = width
                box.top = top
                box.height = height

                j += 1

    def make_content(self):
        for i, competitor in enumerate(self.competitor_order):
            player_boxes = self.player_boxes[i]
            for player in self.player_preds:
                player_box = player_boxes[player]
                text_frame = player_box.text_frame
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = p.add_run()
                pos, score = self.player_preds[player][competitor]
                running_score = self.running_scores[player][competitor]
                run.text = f"{player} - P{pos}\n{running_score} (+{score})"

    def make_titles(self):
        for i, competitor in enumerate(self.competitor_order):
            slide = self.slides[i]
            title_box = make_title_layout(self.prs, slide, 4/25)
            make_title_box(f"{self.base_title}: P{len(self.competitor_order) - i} - {competitor}", title_box)

if __name__ == "__main__":

    prs = Presentation()
    prs.slide_width = 12192 * 1000
    blank_slide_layout = prs.slide_layouts[6]

    spreadsheet_path = Path("C:\Projekty\Coding\Python\F1PredsPPT\F12025 Predictions Tracking.xlsx")
    wb = px.open(spreadsheet_path, data_only=True)

    dc_reader = ChampionshipReader2025(wb["DriverPredictions"])
    dc_reader.gather_data(number_of_competitors=20)
    dc_data = dc_reader.format_to_slide()

    driver_order = dc_data[0]

    player_preds = dc_data[1]

    running_scores = dc_data[2]

    base_title = "Driver's Championship"

    drivers_slide = ChampionshipSlides(
        prs,
        driver_order,
        player_preds,
        running_scores,
        base_title)

    drivers_slide.make_slide()

    prs.save("examples/driversSlideExample.pptx")
