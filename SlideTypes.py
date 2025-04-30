from pptx import Presentation, presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.slide import Slide
from pptx.shapes.autoshape import Shape
from pptx.enum.text import PP_ALIGN

import constants2024 as con
from constants2024 import NUMBER_OF_DRIVERS, NUMBER_OF_PLAYERS, NUMBER_OF_TEAMS
from make_ppt import make_title_box, make_title_layout

class IntroSlide():

    def __init__(
            self,
            prs: presentation.Presentation,
            title: str,
            summary: str,
            player_scores: dict[str, int]
        ):
        self.prs = prs
        blank_slide_layout = prs.slide_layouts[6]
        self.slide = prs.slides.add_slide(blank_slide_layout)

        self.title = title
        self.summary = summary
        self.player_scores = player_scores

        self.title_box = make_title_layout(prs, self.slide, 4/25)
        make_title_box(title, self.title_box)

    def make_boxes(self):
        ph = Inches(1.0)

        shape_tree = self.slide.shapes

        self.summary_box = shape_tree.add_shape(
            MSO_SHAPE.RECTANGLE, ph, ph, ph, ph
        )
        self.player_boxes = {player: shape_tree.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph
        ) for player in self.player_scores}

    def make_layout(self):
        width, height = self.prs.slide_width, self.prs.slide_height

        self.summary_box.left = int(2/25 * width)
        self.summary_box.width = int(12/25 * width)
        self.summary_box.top = int(5/25 * height)
        self.summary_box.height = int(19/25 * height)

        player_boxes = list(self.player_boxes.items())

        for i in range(len(player_boxes)):
            player, box = player_boxes[i]

            player_rgb = con.PLAYERS_TO_COLOURS[player]
            outline_rgb = (int(0.8 * i) for i in player_rgb)
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*player_rgb)
            box.line.color.rgb = RGBColor(*outline_rgb)

            box.left = int(17/25 * width)
            box.width = int(4/25 * width)
            box.top = int(i * 5/50 * height + 5/25 * height)
            box.height = int(2/25 * height)

    def make_content(self):
        sb_tf = self.summary_box.text_frame
        sb_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        run = sb_tf.paragraphs[0].add_run()
        run.text = self.summary

        player_boxes = list(self.player_boxes.items())

        for i in range(len(player_boxes)):
            player, box = player_boxes[i]

            box_tf = box.text_frame
            run = box_tf.paragraphs[0].add_run()
            run.text = f"{player} - {self.player_scores[player]}"

    def make_slide(self):
        self.make_boxes()
        self.make_layout()
        self.make_content()

class H2HSlide():
    '''Displays the'''

    def __init__(
        self,
        prs: presentation.Presentation,
        drivers: tuple[str, str],
        score: tuple[int, int],
        player_preds: dict[str, str]
        ):

        blank_slide_layout = prs.slide_layouts[6]
        self.slide = prs.slides.add_slide(blank_slide_layout)

        self.driver1, self.driver2 = drivers
        self.h2h_score = score

        self.score_left = self.score_right = 0

        self.player_preds = player_preds

        self.prs = prs

    def make_boxes(self):
        '''
        Makes all the boxes on the H2H slide, which are:
        - Box for Driver 1
        - Box for Driver 2
        - The H2H Score box
        - The squares markers representing players.
        '''

        #Placeholder value purely for initiation.
        ph = Inches(1.0)

        shape_tree = self.slide.shapes

        self.driver1_box = shape_tree.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                ph, ph, ph, ph)
        self.driver2_box = shape_tree.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                                ph, ph, ph, ph)
        self.score_box = shape_tree.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                              ph, ph, ph, ph)

        self.player_boxes = {player: shape_tree.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph)
            for player in self.player_preds}

    def make_layout(self):
        '''Places all boxes in right places and makes style.'''

        # The driver boxes share all style and vertical information.
        for driver_box in (self.driver1_box, self.driver2_box):
            driver_box.width = driver_box.height = int(self.prs.slide_height * 10/25)
            driver_box.fill.solid()
            driver_box.fill.fore_color.rgb = RGBColor(1,1,1)

        # 3/25th seems visually correct. We place the driver2 box at 22/25th
        # on its right edge.
        self.driver1_box.left = int(3 / 25 * self.prs.slide_width)
        self.driver2_box.left = self.prs.slide_width - self.driver1_box.width - self.driver1_box.left

        # The Top of the boxes is 1/25th padded to the 5/25th that the title
        # box takes up.
        self.driver1_box.top = self.driver2_box.top = int((5+1) / 25 * self.prs.slide_height)

        # All players share the same size, and have their colour defined in
        # the constants. We make the boxes outline an 80% darker (with rounding)
        # to stand out against the slide background.
        for player in self.player_boxes:
            shape = self.player_boxes[player]

            # Colour stuff
            player_rgb = con.PLAYERS_TO_COLOURS[player]
            outline_rgb = (int(0.8 * colour) for colour in player_rgb)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*player_rgb)
            shape.line.color.rgb = RGBColor(*outline_rgb)

            # Again padding by 1/25th, but this relative to the driver boxes.
            shape.top = int(self.prs.slide_height * (1) / 25) + self.driver1_box.top + self.driver1_box.height
            shape.width = shape.height = int(self.prs.slide_height * 1/25)

        # The score box is centred relative to the driver 1 and driver 2 boxes.
        self.score_box.width = int(4 / 25 * self.prs.slide_width)
        self.score_box.left = int((self.driver2_box.left + self.driver1_box.left + self.driver1_box.width) / 2 - self.score_box.width / 2)
        self.score_box.height = int(self.prs.slide_height * 3 / 25)
        self.score_box.top = self.driver1_box.top + int((self.driver1_box.height - self.score_box.height) / 2)

        # We align the player boxes underneath the driver boxes in a line. This
        # is currently on the edge when all 7 people voted for 1 driver. Could
        # potentially move to 2 lines if more people join in the fun.
        left_shapes, right_shapes = [], []

        for player in self.player_preds:
            player_pred = self.player_preds[player]
            if player_pred == self.driver1:
                left_shapes.append(self.player_boxes[player])
            elif player_pred == self.driver2:
                right_shapes.append(self.player_boxes[player])

        # We want the boxes to be spaced out in a way where they are evenly
        # spaced out between the left and right parts of the driver boxes.
        # The boxes should be centred about the driver box of their choice
        # as a group.
        left_lims = (self.driver1_box.left, self.driver1_box.left + self.driver1_box.width - shape.width)
        right_lims = (self.driver2_box.left, self.driver2_box.left + self.driver2_box.width - shape.width)

        left_divs = len(left_shapes) + 1
        right_divs = len(right_shapes) + 1

        left_unit = int((left_lims[1] - left_lims[0]) / left_divs)

        for i in range(len(left_shapes)):
            shape = left_shapes[i]
            shape.left = left_lims[0] + int((i+1) * left_unit)

        right_unit = int((right_lims[1] - right_lims[0]) / right_divs)

        for i in range(len(right_shapes)):
            shape = right_shapes[i]
            shape.left = right_lims[0] + int((i+1) * right_unit)

    def make_content(self):

        # Add Driver text
        for driver, driver_box in [(self.driver1, self.driver1_box), (self.driver2, self.driver2_box)]:
            text_box = driver_box.text_frame
            p = text_box.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = driver

        # Add Score text.
        text_box = self.score_box.text_frame
        p = text_box.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f"{self.h2h_score[0]}-{self.h2h_score[1]}"

    def make_slide(self):
        '''Calls all the steps to make the H2H slide in order.'''
        self.make_boxes()
        self.make_layout()
        self.make_content()

class RoundUpSlide():

    def __init__(self,
                 prs: presentation.Presentation,
                 player_scores: dict[str, int],
                 title: str):
        self.player_scores = player_scores

        blank_slide_layout = prs.slide_layouts[6]
        self.slide = prs.slides.add_slide(blank_slide_layout)

        self.prs = prs
        self.title = title

        self.make_slide()

    def make_boxes(self):
        #Placeholder value purely for initiation.
        ph = Inches(1.0)

        shape_tree = self.slide.shapes

        self.player_boxes = []
        for player in self.player_scores:
            self.player_boxes.append(
                (self.player_scores[player],
                player,
                shape_tree.add_shape
                    (
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    ph, ph, ph, ph)
                    )
                )
        self.player_boxes.sort(key = lambda x: x[0], reverse=True)

    def make_layout(self):

        width, height = self.prs.slide_width, self.prs.slide_height

        w_prop_box = 3/25
        w_prop_pad = 0
        left_lim, right_lim = (w_prop_pad * width, (1 - w_prop_box - w_prop_pad) * width)

        n_bot = len(self.player_boxes) // 2
        n_top = len(self.player_boxes) - n_bot

        bottom_ls = [int(left_lim + (right_lim - left_lim) * (i+1) / (n_bot+1)) for i in range(n_bot)]
        top_ls = [int(left_lim + (right_lim - left_lim) * (i+1) / (n_top+1)) for i in range(n_top)]

        top_ls.extend(bottom_ls)
        lefts = top_ls

        for i in range(len(lefts)):
            self.player_boxes[i][2].left = lefts[i]
            self.player_boxes[i][2].width = int(w_prop_box * width)
            self.player_boxes[i][2].height = int(1/2 * w_prop_box * width)
            if i > len(lefts) // 2:
                self.player_boxes[i][2].top = int(0.5 * self.player_boxes[i][2].height) + int(1/2 * height)
            else:
                self.player_boxes[i][2].top = int(- 1.5 * self.player_boxes[i][2].height) + int(1/2 * height)

    def make_content(self):
        for score, player, box in self.player_boxes:
            player_rgb = con.PLAYERS_TO_COLOURS[player]
            outline_rgb = (int(0.8 * colour) for colour in player_rgb)
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(*player_rgb)
            box.line.color.rgb = RGBColor(*outline_rgb)

            text_box = box.text_frame
            p = text_box.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = f"{player}\n{score}"

    def make_title(self):
        title_box = make_title_layout(self.prs, self.slide, 4/25)
        make_title_box(self.title, title_box)

    def make_slide(self):
        self.make_boxes()
        self.make_layout()
        self.make_content()
        self.make_title()

class CCSlides():

    ROWS = 4

    def __init__(
        self,
        prs: presentation.Presentation,
        real_order: dict[str, int],
        player_orders: dict[str, dict[str, int]],
        player_scores: dict[str, int]
        ):

        self.real_order = real_order
        self.player_orders = player_orders
        self.player_scores = player_scores

        self.prs = prs
        blank_slide_layout = prs.slide_layouts[6]
        self.slides = {name: prs.slides.add_slide(blank_slide_layout)
                       for name in player_orders}

    def make_boxes(self):
        ph = Inches(1.0)

        self.result_roundup = {}
        self.score_roundup = {}

        for name, slide in self.slides.items():
            shape_tree = slide.shapes

            self.result_roundup[name] = shape_tree.add_table(
                    NUMBER_OF_TEAMS+1, CCSlides.ROWS, ph, ph, int(3/25 * self.prs.slide_width), ph)

            self.score_roundup[name] = shape_tree.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    ph, ph, ph, ph)

            title_text = f"Constructor's - {name}"
            title_box = make_title_layout(self.prs, slide, 4/25)
            make_title_box(title_text, title_box)

    def make_layout(self):

        width, height = self.prs.slide_width, self.prs.slide_height

        for name in self.result_roundup:
            result_roundup = self.result_roundup[name].table
            score_roundup = self.score_roundup[name]

            pad_prop = 3/25
            result_roundup.left = int(pad_prop * width)
            score_roundup.width = int(3/25 * width)
            score_roundup.left = int((1 - pad_prop) * width) - score_roundup.width
            # result_roundup.width =
            table_width = score_roundup.left - result_roundup.left - int(pad_prop * width)

            for col in result_roundup.columns:
                col.width = int(table_width / CCSlides.ROWS)

            h_pad_prop = 5/25

            result_roundup = self.result_roundup[name]

            result_roundup.top = int(h_pad_prop * height)
            result_roundup.height = int((1 - h_pad_prop) * height) - result_roundup.top

            score_roundup.height = int(4/25 * height)
            score_roundup.top = result_roundup.top + int((result_roundup.height  - score_roundup.height)/ 2)

    def make_content(self):
        real_order = sorted([team for team in self.real_order], key= lambda x: self.real_order[x])

        for name, slide in self.slides.items():
            result_roundup = self.result_roundup[name].table
            score_roundup = self.score_roundup[name]

            player_order = sorted([team for team in self.real_order], key= lambda x: self.player_orders[name][x])

            diff_squares = [(self.real_order[team_name] - self.player_orders[name][team_name])**2 for team_name in player_order]
            diff_square = sum(diff_squares)
            for i in range(len(player_order)):
                to_input = (i + 1, real_order[i], player_order[i], diff_squares[i])
                for j in range(CCSlides.ROWS):
                    cell = result_roundup.cell(i + 1, j)
                    cell.text = str(to_input[j])

            to_input = ("Position", "Team", "Predicted Team", "Diff Squared")
            for j in range(CCSlides.ROWS):
                cell = result_roundup.cell(0, j)
                cell.text = str(to_input[j])

            text_frame = score_roundup.text_frame
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()

            run.text = f"Diffs = {diff_square}\nScore:\n{self.player_scores[name]}"

    def make_slide(self):
        self.make_boxes()
        self.make_layout()
        self.make_content()

class DCSlides():

    ROWS = 4

    def __init__(
        self,
        prs: presentation.Presentation,
        real_order: dict[str, int],
        player_orders: dict[str, dict[str, int]],
        player_scores: dict[str, int]
        ):

        self.real_order = real_order
        self.player_orders = player_orders
        self.player_scores = player_scores

        self.prs = prs
        blank_slide_layout = prs.slide_layouts[6]
        self.slides = {name: prs.slides.add_slide(blank_slide_layout)
                       for name in player_orders}

    def make_boxes(self):
        ph = Inches(1.0)

        self.result_roundup = {}
        self.score_roundup = {}

        for name, slide in self.slides.items():
            shape_tree = slide.shapes

            self.result_roundup[name] = shape_tree.add_table(
                    NUMBER_OF_DRIVERS+1, CCSlides.ROWS, ph, ph,
                    int(12/25 * self.prs.slide_width),
                    int(15/25 * self.prs.slide_height))

            self.score_roundup[name] = shape_tree.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    ph, ph, ph, ph)

            title_text = f"Driver's - {name}"
            title_box = make_title_layout(self.prs, slide, 3/25)
            make_title_box(title_text, title_box)

    def make_layout(self):

        width, height = self.prs.slide_width, self.prs.slide_height

        for name, score_box in self.score_roundup.items():
            score_box.left = int(17/25 * width)
            score_box.width = int(4/25 * width)
            score_box.top = int(12/25 * height)
            score_box.height = int(3/25 * height)
            player_rgb = con.PLAYERS_TO_COLOURS[name]
            outline_rgb = (int(0.8 * colour) for colour in player_rgb)
            score_box.fill.solid()
            score_box.fill.rgb = RGBColor(*player_rgb)
            score_box.line.rgb = RGBColor(*outline_rgb)

    def make_content(self):
        real_order = sorted([driver for driver in self.real_order], key= lambda x: self.real_order[x])

        for name, slide in self.slides.items():
            result_roundup = self.result_roundup[name].table
            score_roundup = self.score_roundup[name]

            player_order = sorted([driver for driver in self.real_order], key= lambda x: self.player_orders[name][x])

            diff_squares = [(self.real_order[team_name] - self.player_orders[name][team_name])**2 for team_name in player_order]
            diff_square = sum(diff_squares)
            for i in range(len(player_order)):
                to_input = (i + 1, real_order[i], player_order[i], diff_squares[i])
                for j in range(DCSlides.ROWS):
                    cell = result_roundup.cell(i + 1, j)
                    text_frame = cell.text_frame
                    run = text_frame.paragraphs[0].add_run()
                    run.text = str(to_input[j])
                    run.font.size = Pt(12)
                    cell.margin_bottom = Pt(3)
                    cell.margin_top = Pt(3)

            to_input = ("Position", "Actual Driver", "Predicted Driver", "Diff Squared")
            for j in range(DCSlides.ROWS):
                cell = result_roundup.cell(0, j)
                text_frame = cell.text_frame
                run = text_frame.paragraphs[0].add_run()
                run.text = str(to_input[j])
                run.font.size = Pt(12)
                cell.margin_bottom = Pt(3)
                cell.margin_top = Pt(3)

            text_frame = score_roundup.text_frame
            p = text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()

            run.text = f"Diffs = {diff_square}\nScore:\n{self.player_scores[name]}"

    def make_slide(self):
        self.make_boxes()
        self.make_layout()
        self.make_content()

# class First6RacesSlide():

#     def __init__(self)

class TrueFalseSlide():
    COLS = 3 + NUMBER_OF_PLAYERS

    def __init__(
            self,
            prs: presentation.Presentation,
            driver_cond: dict[str, tuple[int, bool]],
            player_to_truths: dict[str, list[str]],
            player_scores: dict[str, int],
            show_scores: bool = False
        ):

        self.prs = prs
        blank_slide_layout = prs.slide_layouts[6]
        self.slide = prs.slides.add_slide(blank_slide_layout)

        self.driver_cond = driver_cond
        self.player_to_truths = player_to_truths
        self.player_scores = player_scores

        self.show_scores = show_scores
        self.rows = NUMBER_OF_DRIVERS + 1
        if self.show_scores:
            self.rows += 1

    def make_boxes(self):
        ph = Inches(0.5)

        shape_tree = self.slide.shapes
        self.driver_conds = shape_tree.add_table(
            self.rows, TrueFalseSlide.COLS, ph, ph, int(23/25 * self.prs.slide_width), int(10/25 * self.prs.slide_height))
        # self.player_circles = {player: [shape_tree.add_shape(
        #                                 MSO_SHAPE.ROUNDED_RECTANGLE,
        #                                 ph,
        #                                 ph,
        #                                 ph,
        #                                 int(self.prs.slide_height  * 12 / (25 * (NUMBER_OF_DRIVERS+1)))) for i in range(len(drivers))]
        #                         for player, drivers in self.player_to_truths.items()}
        # self.score_boxes = {player: shape_tree.add_shape(
        #                         MSO_SHAPE.ROUNDED_RECTANGLE,
        #                         ph, ph, ph, ph) for player in self.player_to_truths}

    def make_layout(self):
        width, height = self.prs.slide_width, self.prs.slide_height

        self.driver_conds.top = int(5/25 * height)
        self.driver_conds.left = int(1/25 * width)

        # for player, circles in self.player_circles.items():
        #     for circle in circles:
        #         player_rgb = con.PLAYERS_TO_COLOURS[player]
        #         outline_rgb = (int(0.8 * colour) for colour in player_rgb)
        #         circle.fill.solid()
        #         circle.fill.fore_color.rgb = RGBColor(*player_rgb)
        #         circle.line.color.rgb = RGBColor(*outline_rgb)

    def make_content(self):

        driver_names = list(self.driver_cond.keys())

        # driver_to_height = {}
        cell_height = self.driver_conds.height / (NUMBER_OF_DRIVERS + 1)
        players = tuple(self.player_to_truths.keys())

        for i in range(1, NUMBER_OF_DRIVERS + 1):
            driver = driver_names[i - 1]
            data = (driver, *self.driver_cond[driver], *(" " for i in range(NUMBER_OF_PLAYERS)))
            row = self.driver_conds.table.rows[i - 1]
            row.height = int(cell_height)

            driver_true = self.driver_cond[driver][1]
            for j in range(TrueFalseSlide.COLS):
                cell = self.driver_conds.table.cell(i, j)
                text_frame = cell.text_frame
                run = text_frame.paragraphs[0].add_run()
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run.text = f"{data[j]}"
                run.font.size = Pt(10)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.margin_bottom = cell.margin_top = Pt(3)
                if j > 2:
                    driver_predicted = driver in self.player_to_truths[players[j - 3]]
                    if driver_predicted:
                        cell.fill.solid()
                        player_rgb = con.PLAYERS_TO_COLOURS[players[j - 3]]
                        cell.fill.fore_color.rgb = RGBColor(*player_rgb)

                    if self.show_scores:
                        # print(data, driver_true, driver_predicted)
                        if driver_predicted and driver_true:
                            run.text = "+5"
                        elif not driver_predicted and not driver_true:
                            run.text = "+0"
                        else:
                            run.text = "-3"

        if self.show_scores:
            player_scores = tuple(self.player_scores.values())
            row_data = (" ", " ", " ", *player_scores)
            # print(player_scores)
            for j in range(TrueFalseSlide.COLS):
                cell = self.driver_conds.table.cell(NUMBER_OF_DRIVERS + 1, j)
                text_frame = cell.text_frame
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.CENTER
                run = text_frame.paragraphs[0].add_run()
                run.text = f"{row_data[j]}"
                run.font.size = Pt(10)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(255, 255, 255)
                cell.margin_bottom = cell.margin_top = Pt(3)

        title_data = ("Driver", "Count", "True/False", *players)
        for j in range(TrueFalseSlide.COLS):
            cell = self.driver_conds.table.cell(0, j)
            text_frame = cell.text_frame
            run = text_frame.paragraphs[0].add_run()
            run.text = f"{title_data[j]}"
            run.font.size = Pt(10)
            # driver_to_height[driver] = i * (cell.margin_bottom + cell.margin_top + cell_height)  + self.driver_conds.top

        # players = list(self.player_circles.keys())
        # circle_width = int((10/25 * self.prs.slide_width) / NUMBER_OF_PLAYERS)
        # circle_delta = int((8/25 * self.prs.slide_width) / NUMBER_OF_PLAYERS)

        # for i in range(len(players)):
        #     player = players[i]
        #     player_conds = self.player_to_truths[player]
        #     player_circles = self.player_circles[player]
        #     for j in range(len(player_conds)):
        #         driver, circle = player_conds[j], player_circles[j]
        #         circle.top = int(driver_to_height[driver])
        #         circle.height = int(cell_height)
        #         circle.width = int(cell_height)
        #         circle.left = int(15/25 * self.prs.slide_width + i * circle_delta)

    def make_slide(self):
        self.make_boxes()
        self.make_layout()
        self.make_content()

class SuperlativeSlide():
    def __init__(
            self,
            prs: presentation.Presentation,
            driver_rankings: list[tuple[float, str, int]],
            player_preds: dict[str, tuple[str, int]],
            slide_type: str = "Driver"
            ):
        self.prs = prs
        blank_slide_layout = prs.slide_layouts[6]
        self.slide = prs.slides.add_slide(blank_slide_layout)

        self.driver_rankings = driver_rankings
        self.player_preds = player_preds
        self.type = slide_type

    def make_boxes(self):
        ph = Inches(1.0)
        shape_tree = self.slide.shapes
        self.player_boxes = {player: shape_tree.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph
            ) for player in self.player_preds}

    def make_layout(self):
        width, height = self.prs.slide_width, self.prs.slide_height

        left = int(1/25 * width)
        t_width = int(10/25 * width)
        top = int(5/25 * height)
        t_height = int(14/25 * height)

        self.driver_roundup = self.slide.shapes.add_table(
            len(self.driver_rankings) + 1, 3, left, top, t_width, t_height
            )
        players = list(self.player_boxes.keys())

        for i in range(len(players)):
            player = players[i]
            shape = self.player_boxes[player]
            player_rgb = con.PLAYERS_TO_COLOURS[player]
            outline_rgb = (int(0.8 * colour) for colour in player_rgb)
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*player_rgb)
            shape.line.color.rgb = RGBColor(*outline_rgb)

            shape.left = int(16/25 * width)
            shape.width = int(4/25 * width)
            shape.height = int(3/50 * height)
            shape.top = self.driver_roundup.top + int(i * 2/25 * height)

    def make_content(self):

        for player in self.player_boxes:
            shape = self.player_boxes[player]
            text = f"{player}\n{self.player_preds[player][0]} - {self.player_preds[player][1]}"
            tf = shape.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.add_run()
            run.text = text
            run.font.size = Pt(12)

        data = [("Rank", self.type, "Count")]
        data.extend(self.driver_rankings)
        for i in range(len(data)):
            for j in range(len(data[0])):
                cell = self.driver_roundup.table.cell(i, j)
                text_frame = cell.text_frame
                run = text_frame.paragraphs[0].add_run()
                p = text_frame.paragraphs[0]
                p.alignment = PP_ALIGN.LEFT
                run.text = f"{data[i][j]}"
                run.font.size = Pt(12)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(200, 200, 200)
                cell.margin_bottom = cell.margin_top = Pt(5)

    def make_slide(self):
        self.make_boxes()
        self.make_layout()
        self.make_content()


class First6RacesSlides():
    def __init__(
            self,
            prs: presentation.Presentation,
            driver_orders: dict[str, list[int]],
            player_preds: dict[str, list[tuple[str, int]]]):
        self.prs = prs
        blank_slide_layout = prs.slide_layouts[6]
        self.slides = [prs.slides.add_slide(blank_slide_layout) for i in range(6)]

        self.driver_orders = driver_orders
        self.player_preds = player_preds

    def make_boxes(self):
        ph = Inches(1.0)

        self.driver_scores = []
        self.player_boxes = {player: [] for player in self.player_preds}

        for slide in self.slides:
            shape_tree = slide.shapes
            self.driver_scores.append(shape_tree.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph
            ))
            for player in self.player_boxes:
                self.player_boxes[player].append(shape_tree.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE, ph, ph, ph, ph
            ))

    def make_layout(self):
        width, height = self.prs.slide_width, self.prs.slide_height
        players = list(self.player_boxes.keys())

        for i in range(len(self.slides)):
            # slide = self.slides[i]
            driver_roundup = self.driver_scores[i]
            driver_roundup.left = int(1/25 * width)
            driver_roundup.width = int(12/25 * width)
            driver_roundup.top = int(5/25 * height)
            driver_roundup.height = int(18/25 * height)

            for j in range(len(players)):
                player = players[j]
                player_box = self.player_boxes[player][i]
                player_box.left = int(15/25 * width)
                player_box.width = int(5/25 * width)
                player_box.top = int((5/25 + j/10) * height)
                player_box.height = int(2/25 * height)

                player_rgb = con.PLAYERS_TO_COLOURS[player]
                outline_rgb = (int(0.8 * i) for i in player_rgb)
                player_box.fill.solid()
                player_box.fill.fore_color.rgb = RGBColor(*player_rgb)
                player_box.line.color.rgb = RGBColor(*outline_rgb)

    def make_content(self):
        partial_scores = {player: 0 for player in self.player_preds}

        for i in range(len(self.slides)):
            drivers_to_display = []
            for driver in self.driver_orders:
                driver_pos = self.driver_orders[driver][i]
                shifted_pos = abs(driver_pos - i - 1) + 1
                if shifted_pos < 11:
                    drivers_to_display.append((driver_pos, shifted_pos, driver))
            text = ""
            drivers_to_display.sort(key = lambda x: int(x[0]))
            for (pos, spos, driver) in drivers_to_display:
                text = text + f"{pos} - {spos} - {driver}\n"
            tf = self.driver_scores[i].text_frame
            tf.paragraphs[0].align = PP_ALIGN.LEFT
            run = tf.paragraphs[0].add_run()
            run.text = text

            for player, pred in self.player_preds.items():
                player_box = self.player_boxes[player][i]
                driver_pred, score = pred[i]
                partial_scores[player] += score
                tf = player_box.text_frame
                run = tf.paragraphs[0].add_run()
                run.text = f"{player}: {driver_pred} \n +{score}, {partial_scores[player]}"
                run.font.size = Pt(12)

    def make_titles(self):
        for i in range(len(self.slides)):
            slide = self.slides[i]
            title = make_title_layout(self.prs, slide, 4/25)
            make_title_box(f"First 6 Races - Round {i+1}", title)

    def make_slide(self):
        self.make_boxes()
        self.make_layout()
        self.make_content()
        self.make_titles()

class Pick5RacesSlide():
    ROWS = 2 * NUMBER_OF_PLAYERS
    COLUMNS = 5 + 1

    def __init__(
            self,
            prs: presentation.Presentation,
            player_preds: dict[str, list[tuple[str, int]]],
            all_races: set[str],
            running_scores: dict[str, int],
            base_title: str):
        self.prs = prs
        blank_slide_layout = prs.slide_layouts[6]
        # self.blank_slide = prs.slides.add_slide(blank_slide_layout)
        self.slides = [prs.slides.add_slide(blank_slide_layout) for race in con.RACES if race in all_races]

        self.player_preds = player_preds
        self.all_races = all_races
        self.running_scores = running_scores
        self.base_title = base_title

    def make_boxes(self):
        ph = Inches(1.0)

        self.pick5_tables = [slide.shapes.add_table(
                    Pick5RacesSlide.ROWS, Pick5RacesSlide.COLUMNS,
                    int(1/25 * self.prs.slide_width), int(5/25 * self.prs.slide_height),
                    int(18/25 * self.prs.slide_width), int(15/25 * self.prs.slide_height))
                    for slide in self.slides]
        self.player_boxes = [
            {player: slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, ph, ph, ph, ph
            ) for player in self.player_preds}
            for slide in self.slides
        ]

    def make_layout(self):
        width, height = self.prs.slide_width, self.prs.slide_height

        players = list(self.player_preds.keys())
        for i in range(len(self.slides)):
            player_boxes = self.player_boxes[i]
            pick5_table = self.pick5_tables[i].table
            for j in range(len(players)):
                player = players[j]
                player_box = player_boxes[player]

                player_box.left = int(20/25 * width)
                player_box.width = int(3/25 * width)
                player_box.top = int((6/25 + j/10) * height)
                player_box.height = int(2/25 * height)

                player_rgb = con.PLAYERS_TO_COLOURS[player]
                outline_rgb = (int(0.8 * i) for i in player_rgb)
                player_box.fill.solid()
                player_box.fill.fore_color.rgb = RGBColor(*player_rgb)
                player_box.line.color.rgb = RGBColor(*outline_rgb)

                for k in range(Pick5RacesSlide.COLUMNS):
                    for m in (2 * j, 2 * j + 1):
                        cell = pick5_table.cell(m, k)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = RGBColor(*player_rgb)

    def make_content(self):
        to_show = {player: [] for player in self.player_preds}
        running_round_score = {player: 0 for player in self.player_preds}
        player_pred_q = self.player_preds.copy()
        players = list(self.player_preds.keys())

        i = 0
        for race in con.RACES:
            player_boxes = self.player_boxes[i]
            pick5_table = self.pick5_tables[i].table
            if race in self.all_races:
                for player, remaining_preds in player_pred_q.items():
                    if len(remaining_preds) != 0:
                        player_race, score = remaining_preds[0]
                        if player_race == race:
                            to_show[player].append(score)
                            player_pred_q[player] = remaining_preds[1:]
                            running_round_score[player] += score
                            self.running_scores[player] += score

                for j in range(len(players)):
                    player = players[j]
                    nr_row = 2 * j
                    player_and_races = [player] + [race for race, score in self.player_preds[player]]
                    for k in range(Pick5RacesSlide.COLUMNS):
                        # Name, Races Row
                        cell = pick5_table.cell(nr_row, k)
                        tf = cell.text_frame
                        run = tf.paragraphs[0].add_run()
                        run.text = player_and_races[k]
                    score_row = 2 * j + 1
                    pscore_and_scores = [running_round_score[player]] + to_show[player]
                    for k in range(len(pscore_and_scores)):
                        cell = pick5_table.cell(score_row, k)
                        tf = cell.text_frame
                        run = tf.paragraphs[0].add_run()
                        run.text = str(pscore_and_scores[k])
                    player_box = player_boxes[player]
                    tf = player_box.text_frame
                    run = tf.paragraphs[0].add_run()
                    run.text = str(self.running_scores[player])

                i += 1

    def make_titles(self):
        i = 0
        for race in con.RACES:
            slide = self.slides[i]
            if race in self.all_races:
                title_box = make_title_layout(self.prs, slide, 4/25)
                make_title_box(f"{self.base_title} @ {race}", title_box)
                i += 1

    def make_slide(self):
        self.make_titles()
        self.make_boxes()
        self.make_layout()
        self.make_content()

if __name__ == "__main__":

    prs = Presentation()
    prs.slide_width = 12192 * 1000
    blank_slide_layout = prs.slide_layouts[6]

    # print(prs.slide_width / Inches(1.0))

    drivers = ("Max Verstappen", "Sergio Perez")
    score = (23, 1)

    players = {
        "Benedict": "Max Verstappen",
        "Carla": "Max Verstappen",
        "Damian": "Max Verstappen",
        "Jarek": "Max Verstappen",
        "Josh": "Sergio Perez",
        "Kacper": "Max Verstappen",
        "Suley": "Max Verstappen",
    }

    h2h_eg = H2HSlide(prs, drivers, score, players)
    h2h_eg.make_boxes()
    h2h_eg.make_layout()
    h2h_eg.make_content()

    roundup_eg = RoundUpSlide(prs, {player: int(player[:2], 36) for player in players})
    roundup_eg.make_slide()
    # roundup_eg.make_layout()

    tf_slide = TrueFalseSlide(prs, {}, {}, {})
    tf_slide.make_slide()

    prs.save("slideTypeExamples.pptx")