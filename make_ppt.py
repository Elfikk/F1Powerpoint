from pptx import Presentation, presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.slide import Slide
from pptx.shapes.autoshape import Shape
from pptx.enum.text import PP_ALIGN

import constants2024 as con

def make_player_box(shape: Shape, name = "AA", score = "69420"):

    text_frame = shape.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = f'{name}: {score}'

def make_player_layout(
        prs: presentation,
        slide: Slide,
        num_players: int = 7,
        available_h: float = 6/25
        ) -> list[Shape]:
    # Adds player shapes at correct positions.

    k = num_players // 2

    # if numplayers is odd
    if num_players % 2:
        divs = 3 * k + 4
        top_ls = list(range(1, divs, 3))
        bottom_ls = [i + 0.5 for i in range(2, 3 * k + 2, 3)]
    else:
        divs = 3 * k + 1
        top_ls = [i for i in range(1, divs, 3)]
        bottom_ls = top_ls

    delta_h = available_h / 2

    top_h = (1 - available_h) * prs.slide_height
    bottom_h = top_h + delta_h * prs.slide_height

    shape_tree = slide.shapes

    width = 2 / divs * prs.slide_width
    height = delta_h * prs.slide_height

    player_shapes = []

    for i in top_ls:
        left = i / divs * prs.slide_width
        player_shape = shape_tree.add_shape(MSO_SHAPE.RECTANGLE, left, top_h, width, height)
        make_player_box(player_shape)
        player_shapes.append(player_shape)

    for i in bottom_ls:
        left = i / divs * prs.slide_width
        player_shape = shape_tree.add_shape(MSO_SHAPE.RECTANGLE, left, bottom_h, width, height)
        make_player_box(player_shape)
        player_shapes.append(player_shape)

    return player_shapes

def set_player_detail(players: list[tuple[str, int]], player_boxes: list[Shape]):
    for i in range(len(players)):

        player_name, player_score = players[i]
        player_box = player_boxes[i]

        text_frame = player_box.text_frame
        text_frame.clear()
        p = text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = f'{player_name}: {player_score}'

        player_box.fill.solid()
        rgb = con.PLAYERS_TO_COLOURS[player_name]
        outline_rgb = (int(0.8 * colour) for colour in rgb)
        player_box.fill.fore_color.rgb = RGBColor(*rgb)
        player_box.line.color.rgb = RGBColor(*outline_rgb)

def make_title_box(slide_title: str, title_box: Shape):

    text_frame = title_box.text_frame
    p = text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = slide_title

def make_title_layout(prs, slide: Slide, available_h = 5 / 25) -> Shape:
    shapes = slide.shapes
    title_shape = shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, available_h * prs.slide_height)
    return title_shape

def make_round_roundup_content(player_scores: dict[str, int]):
    # For any finished round, take the scores and display them
    # Then display final scores
    pass

def make_constructor_content(
    real_order: dict[str, int],
    player_order: dict[str, int],
    diffs: dict[str, int],
    player_score: int):
    # The constructor content slide consists of
    # 1. The real constructor order on the LHS
    # 2. The player's predicted constructor order on the RHS
    # 3. The differences squared between the two
    # 4. The player's score.
    pass

def make_drivers_content(
    real_order: dict[str, int],
    player_order: dict[str, int],
    diffs: dict[str, int],
    player_score: int):
    # The driver's championship slides consist of
    # 1. The real driver's order spread over two pages.
    # 2. The player's predicted positions next to each driver.
    # 3. The diffs squared.
    # 4. The player's score.
    pass

def make_six_content(
    partial_orders: list[dict[str, int]],
    player_preds: list[dict[str, int]],
    player_scores: dict[str, int]
):
    # This one is a little bit more convoluted.
    pass

def bool_content():
    # This one is a little bit more convoluted.
    pass

def superlative_content_drivers():
    pass

def superlative_content_teams():
    pass

def five_races_content():
    pass

if __name__ == "__main__":

    prs = Presentation()
    # print(prs.slide_width, prs.slide_height)
    prs.slide_width = 12192 * 1000
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    shapes = slide.shapes

    # left = top = width = height = Inches(1.0)

    left = prs.slide_width / 13
    top = 19 / 25 * prs.slide_height
    width = 2 * prs.slide_width / 13
    height = 3 / 25 * prs.slide_height

    player_boxes = make_player_layout(slide, 7)

    players = [
        ("Benedict", 1),
        ("Carla", 2),
        ("Damian", 3),
        ("Jarek", 4),
        ("Josh", 5),
        ("Kacper", 6),
        ("Suley", 7),
    ]

    set_player_detail(players, player_boxes)

    title = make_title_layout(slide)

    make_title_box("Example Title", title)

    prs.save("text.pptx")